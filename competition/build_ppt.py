# -*- coding: utf-8 -*-
"""设计说明 PPT — 中建二局“建证超越”青年设计师创新大赛决赛"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK = RGBColor(0x0D, 0x18, 0x30)
DIM = RGBColor(0x45, 0x56, 0x7A)
BLUE = RGBColor(0x1B, 0x6F, 0xC4)
BLUE2 = RGBColor(0x2F, 0x6D, 0xF0)
CYAN = RGBColor(0x0F, 0xB5, 0xC9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FAINT = RGBColor(0x85, 0x93, 0xAB)
DARK = RGBColor(0x06, 0x09, 0x0F)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid() if fill else sp.fill.background()
    if fill:
        sp.fill.fore_color.rgb = fill
    if line:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, size=18, color=DIM, bold=False, align=PP_ALIGN.LEFT,
         line_spacing=1.3, anchor=MSO_ANCHOR.TOP, font="Microsoft YaHei"):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        if isinstance(para, tuple):
            txt, kw = para
        else:
            txt, kw = para, {}
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(kw.get("size", size))
        r.font.bold = kw.get("bold", bold)
        r.font.color.rgb = kw.get("color", color)
        r.font.name = kw.get("font", font)
    return tb


def bullets(s, x, y, w, h, items, size=16, gap=10, color=DIM, line_spacing=1.35):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for head, body in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = line_spacing
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = "▎" + head
        r.font.size = Pt(size + 2); r.font.bold = True; r.font.color.rgb = BLUE
        r.font.name = "Microsoft YaHei"
        if body:
            p2 = tf.add_paragraph(); p2.line_spacing = line_spacing; p2.space_after = Pt(gap)
            r2 = p2.add_run(); r2.text = body
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = "Microsoft YaHei"
    return tb


def pic_cover(s, path, x, y, w, h):
    """add picture covering the frame (center-crop)"""
    from PIL import Image
    im = Image.open(path)
    ir = im.width / im.height
    fr = w / h
    pic = s.shapes.add_picture(path, x, y, w, h)
    if ir > fr:  # image wider -> crop left/right
        crop = (1 - fr / ir) / 2
        pic.crop_left = crop; pic.crop_right = crop
    else:
        crop = (1 - ir / fr) / 2
        pic.crop_top = crop; pic.crop_bottom = crop
    return pic


def page_header(s, title, kick):
    rect(s, 0, 0, SW, Inches(1.02), fill=WHITE)
    rect(s, Inches(0.55), Inches(0.36), Inches(0.09), Inches(0.34), fill=BLUE2)
    text(s, Inches(0.78), Inches(0.18), Inches(11.2), Inches(0.5),
         [(kick, {"size": 13, "color": BLUE, "bold": True})])
    text(s, Inches(0.78), Inches(0.44), Inches(11.2), Inches(0.55),
         [(title, {"size": 26, "color": INK, "bold": True})])
    text(s, Inches(10.2), Inches(0.3), Inches(2.7), Inches(0.5),
         [("海外业务赛道 · 向鑫", {"size": 12, "color": FAINT})], align=PP_ALIGN.RIGHT)
    rect(s, 0, Inches(1.02), SW, Pt(2), fill=RGBColor(0xE2, 0xE9, 0xF2))


# ============ S1 cover ============
s = slide()
# 片头动画：浅色 banner 视频满幅嵌入（自动播放，见文件末尾 XML 后处理）
movie = s.shapes.add_movie("competition/img/flow-banner-light-v2.mp4", 0, 0, SW, SH,
                           poster_frame_image="competition/img/hero-light-v2.jpg", mime_type="video/mp4")
# brand row
s.shapes.add_picture("competition/img/company.png", Inches(0.7), Inches(0.55), height=Inches(0.5))
text(s, Inches(0.7), Inches(2.2), Inches(11.9), Inches(0.5),
     [("中建二局“建证·超越”青年设计师创新大赛 · 决赛", {"size": 18, "color": BLUE, "bold": True})])
text(s, Inches(0.7), Inches(2.75), Inches(11.9), Inches(2.2),
     [("海外工程投标及商业调研", {"size": 44, "color": INK, "bold": True}),
      ("全流程 AI 智能 Agent 作业系统", {"size": 44, "color": BLUE, "bold": True})], line_spacing=1.25)
text(s, Inches(0.7), Inches(5.15), Inches(11.9), Inches(0.9),
     [("把海外投标这盘硬活，一次长程任务从头到尾跑完", {"size": 20, "color": DIM})])
rect(s, Inches(0.7), Inches(6.05), Inches(11.9), Pt(1.6), fill=RGBColor(0xC9, 0xD6, 0xE8))
text(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.8),
     [("申报类型：海外业务赛道　｜　参赛人员：向鑫　｜　申报单位：博茨公司", {"size": 16, "color": DIM})])

# ============ S2 设计总说明 ============
s = slide()
page_header(s, "设计总说明", "DESIGN STATEMENT · 1000 字以内")
statement = (
    "「海外工程投标及商业调研全流程 AI 智能 Agent 作业系统」（Agent Pi DSH）是面向工程企业的垂直智能体作业系统。"
    "本人在作品中担负系统架构设计、核心流程与提示工程、投标领域技能与知识库体系搭建、以及真实项目验证的全部工作。"
)
p2 = (
    "海外标段投标长期面临三大痛点：标书文件浩瀚、人工通读易遗漏；数千条 BOQ 组价依赖个人经验、不可校验不可复用；"
    "通用 AI 助手幻觉频发、长任务断片跑题，无法交付高证据化的正式成果。创意构思围绕一个问题展开："
    "普通工程企业的普通电脑，能否把最难的投标任务跑到底。方案做三个反向取舍：其一，知识库不采用市场流行的 RAG 向量检索"
    "——其对普通用户硬件要求过高，且“相似度猜测”不满足证据化场景，改为按文档自身章节切条 + BM25 检索，零重型依赖、命中即定位条款；"
    "其二，依托 DeepSeek V4 百万 tokens 上下文，整套标书、规范、BOQ 与中间稿全程驻留，任务不失忆、目标不偏离；"
    "其三，发动机从 Claude Agent SDK + Pi 双架构果断切换到 DeepSeek Harness 架构，实测最难的 BOQ 逐条分析效果显著更优。"
)
p3 = (
    "系统以一次长程任务跑完投标全流程：标书全量解析并入库 → 数千条 BOQ 逐项引用规范界定工作范围 → "
    "结合企业数据与网络验证的工法、工效、资源实地价格五步推导 → 资源汇总与成本推定 → 按项目特征施工推演形成策划稿 → "
    "照企业模板编制正式投标文档。证据门禁杜绝幻觉，出处芯片逐条回指源文件条款，正式成果落盘可审计；"
    "中标后全部基础资料直接服务实施阶段成本策划。"
)
p4 = (
    "系统已在真实海外项目作业中验证：一次任务生成 EB Cloete 立交 110m 钢拱吊装可交互三维仿真（含 FEM 校验的受力可视化与计算书）、"
    "咨询级版式的投资尽调报告与矿区交通走廊地图。商业化落地按“最小投资私有化”设计：约 3 万元一台工作站即可起步，"
    "软件全部开源零许可费，数据不出企业内网；规划接入微信 / 飞书 / 钉钉的常驻业务中枢，实现企业信息化向 AI 智能化的跃迁。"
)
text(s, Inches(0.78), Inches(1.35), Inches(11.8), Inches(5.7),
     [(statement, {"size": 15.5, "color": INK, "bold": False}),
      (p2, {"size": 15.5, "color": DIM}),
      (p3, {"size": 15.5, "color": DIM}),
      (p4, {"size": 15.5, "color": DIM})], line_spacing=1.5)

# ============ S3 问题提出 ============
s = slide()
page_header(s, "问题提出：海外投标的三大痛点", "PROBLEM")
bullets(s, Inches(0.78), Inches(1.5), Inches(11.8), Inches(5.5), [
    ("标书浩瀚，人工读不完", "海外标段招标文件动辄数十份、上千页，规范、FIDIC 条款与特别条款修订相互交织。人工通读周期长、遗漏多，关键条款理解偏差直接导致废标风险。"),
    ("BOQ 组价靠经验拍脑袋", "数千条清单逐条组价：工作范围界定、工法工效选择、资源实地价格取定，长期依赖少数骨干个人经验，不可校验、不可复用，人走经验走。"),
    ("通用 AI 助手不堪大用", "通用办公助手几十轮对话就断片跑题、凭模型记忆编造事实（幻觉）、长文档读不动、成果只是聊天记录——在高证据化的投标场景里无法交付正式成果。"),
    ("破题思路", "把智能体放进真实项目作业：打开工作区直接下任务；复杂活由内核原生拆给并行工人，证据门禁杜绝幻觉，正式成果带出处芯片落盘，中标后全部基础资料直接服务实施阶段。"),
], size=16, gap=16)

# ============ S4 方案架构 ============
s = slide()
page_header(s, "总体方案：发动机是 DeepSeek Harness，工作台是 Agent Pi", "SYSTEM OVERVIEW")
steps = ["打开项目下任务", "工作台\n投标/实施/投资·知识库", "DeepSeek Harness\n内核·工具/会话/权限", "并行工人\nsubagent / workflow", "Official Outputs\n统一成果树·总报告"]
x = Inches(0.55)
w = Inches(2.32)
for i, t in enumerate(steps):
    c = rect(s, x, Inches(1.7), w, Inches(1.15), fill=(BLUE2 if i == 2 else RGBColor(0xEE, 0xF4, 0xFB)))
    text(s, x, Inches(1.82), w, Inches(0.95), [(t, {"size": 13.5, "color": (WHITE if i == 2 else INK), "bold": True})],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 4:
        text(s, x + w, Inches(1.95), Inches(0.16), Inches(0.6), [("→", {"size": 20, "color": BLUE, "bold": True})], align=PP_ALIGN.CENTER)
    x += w + Inches(0.16)
bullets(s, Inches(0.78), Inches(3.3), Inches(11.8), Inches(3.8), [
    ("证据门禁", "项目特征缺口不能用模型记忆填：要么找到出处，要么由人尽调后授权放行——基本杜绝幻觉。"),
    ("本地知识库", "规范、FIDIC 条款按文档自身章节切条入库，BM25 检索命中即定位 Clause；用户模板复刻企业版式；.apkb 传递包随项目流转。"),
    ("出处芯片", "引用是出处不是摘抄：芯片只显示源文件、页或行、题目或段落，点击即达原文。"),
    ("成果落盘", "统一成果树写回项目目录，崩溃重启不重做；正式稿上的人工修改留在原处。"),
], size=15, gap=10)

# ============ S5 核心流程 ============
s = slide()
page_header(s, "核心流程：一次长程任务，从标书到标稿", "WORKFLOW")
flow = [
    ("01 标书全量解析，规范入库", "一次任务读完所有标书文件；规范、FIDIC 条款入库，与特别条款修订逐条对照总结。"),
    ("02 数千条 BOQ 逐项界定范围", "逐条引用规范、特别条款、FIDIC 修订界定工作范围，每条都有出处芯片。"),
    ("03 五步推导算出每条单价", "结合企业数据与网络验证的工法、工效、资源实时实地价格，详尽五步推导。"),
    ("04 资源汇总与成本推定", "汇总资源清单与成本推定，带公式的组价测算表可直接改。"),
    ("05 按项目特征施工推演", "建立一整套可执行的施工策划稿，不是空话套话。"),
    ("06 照企业模板编制标稿", "用户模板复刻版式、大纲与深度，项目事实永远来自本项目资料。"),
]
x0, y0 = Inches(0.55), Inches(1.5)
for i, (h, b) in enumerate(flow):
    cx = x0 + (i % 3) * Inches(4.2)
    cy = y0 + (i // 3) * Inches(2.55)
    rect(s, cx, cy, Inches(4.0), Inches(2.3), fill=RGBColor(0xEE, 0xF4, 0xFB))
    text(s, cx + Inches(0.25), cy + Inches(0.22), Inches(3.5), Inches(0.6), [(h, {"size": 15.5, "color": BLUE, "bold": True})])
    text(s, cx + Inches(0.25), cy + Inches(0.85), Inches(3.5), Inches(1.3), [(b, {"size": 13, "color": DIM})], line_spacing=1.35)

# ============ S6 架构决策 ============
s = slide()
page_header(s, "三个深思熟虑的架构决策", "DESIGN DECISIONS")
bullets(s, Inches(0.78), Inches(1.5), Inches(11.8), Inches(5.5), [
    ("知识库不堆 RAG：普通电脑就跑得动", "RAG 要嵌入模型、向量库、GPU 与 Docker 编排——太重、用不起；向量检索是「相似度猜测」，检索到 ≠ 证据对。改为按文档自身章/节/条切分 + MiniSearch BM25：零重型依赖，每条命中精确定位到 Clause，出处当场可核。"),
    ("百万上下文：任务全程不失忆", "DeepSeek V4 全系 1M tokens（约 800 页 PDF）——整套标书、规范、数千条 BOQ 与全部中间稿全程驻留上下文。模型不靠检索「猜」，而是在完整事实面前工作：任务始终围绕目标，不因截断而偏离。"),
    ("换到 DSH 架构：实测优于双架构", "最难的 BOQ 逐条分析，DeepSeek Harness 架构运行效果明显优于 Claude Agent SDK + Pi 双架构——内核原生并行工人、持久 Bash 不再卡顿、长任务崩溃只救未完工部分。3.0 果断换芯：回合更短、token 更省、目标更稳。"),
], size=16, gap=18)

# ============ S7 效果表现 ============
s = slide()
page_header(s, "效果表现：来自真实项目作业，不是演示摆拍", "REAL OUTPUTS")
shots = [
    ("competition/img/kb-workbench.png", "知识库工作台：规范入库，BOQ 成果树逐条落盘"),
    ("competition/img/report-equity.png", "投资尽调报告：咨询级版式一次成稿"),
    ("competition/img/report-stages.png", "市场尽调图表：开发阶段全景，来源逐条标注"),
    ("competition/img/report-map.png", "矿区交通走廊地图：原创改绘示意图"),
]
pos = [(Inches(0.55), Inches(1.4)), (Inches(6.95), Inches(1.4)), (Inches(0.55), Inches(4.35)), (Inches(6.95), Inches(4.35))]
for (px, py), (img, cap) in zip(pos, shots):
    pic_cover(s, img, px, py, Inches(5.85), Inches(2.55))
    text(s, px, py + Inches(2.6), Inches(5.85), Inches(0.4), [(cap, {"size": 12.5, "color": DIM})])

# ============ S8 落地推广 ============
s = slide()
page_header(s, "落地与推广：从信息化到智能化的跃迁，可以从一台服务器开始", "DEPLOYMENT")
rows = [
    ("起步型（10–20 人）", "一台高配工作站：24GB 主流显卡 + 64GB 内存 + 2TB 固态硬盘，接入内网即可使用", "约 3 万元"),
    ("标准型（30–80 人）", "双显卡服务器 + 128GB 内存 + 万兆内网，全部业务流程常驻运行", "约 7–10 万元"),
    ("软件", "Agent Pi DSH 服务端、Qwen3.8-27B 本地多模态模型、MinerU、Docker——全部开源", "0 元"),
    ("网络", "内网即可运行，数据不出内网；推送微信/飞书/钉钉只需服务器出站联网", "0 元"),
]
y = Inches(1.5)
for name, desc, cost in rows:
    rect(s, Inches(0.55), y, Inches(12.25), Inches(0.92), fill=RGBColor(0xEE, 0xF4, 0xFB))
    text(s, Inches(0.8), y + Inches(0.12), Inches(2.6), Inches(0.7), [(name, {"size": 15, "color": INK, "bold": True})], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(3.5), y + Inches(0.12), Inches(7.4), Inches(0.7), [(desc, {"size": 13, "color": DIM})], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(10.9), y + Inches(0.12), Inches(1.8), Inches(0.7), [(cost, {"size": 16, "color": BLUE, "bold": True})], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.04)
bullets(s, Inches(0.78), Inches(5.9), Inches(11.8), Inches(1.4), [
    ("推广路线", "PDF 图纸工程量核算（进行中）→ 企业私有化部署：服务端 Web 多用户 + 微信/飞书/钉钉/QQ 接入 + OA 流程插件 → 建筑概念图 / 3D 建模 / BIM 信息化。"),
], size=14, gap=6)

# ============ S9 创作过程资料 ============
s = slide()
page_header(s, "创作过程资料", "PROCESS")
bullets(s, Inches(0.78), Inches(1.5), Inches(11.8), Inches(1.8), [
    ("真实项目验证", "系统已在海外真实项目作业中验证：一次任务生成 EB Cloete 立交 S0153-2 钢拱吊装可交互三维仿真——110m 跨径、分节吊装动画、应力热力图（经 MS FEM 校验）、参数与图纸对照、完整计算书。"),
    ("产品官网与文档", "作品已形成完整产品形态：官方网站（www.agent-pi.app）、中英双语使用文档、34 个出厂领域技能、插件市场与知识库体系。"),
], size=14, gap=8)
pic_cover(s, "competition/img/kb-workbench.png", Inches(0.55), Inches(3.6), Inches(7.9), Inches(3.3))
pic_cover(s, "competition/img/screenshot-market.png", Inches(8.75), Inches(3.6), Inches(4.05), Inches(3.3))

# ============ S10 结束页 ============
s = slide()
pic_cover(s, "competition/img/hero-light-v2.jpg", 0, 0, SW, SH)
text(s, Inches(0.7), Inches(2.9), Inches(11.9), Inches(1.4),
     [("谢谢聆听 · 恳请评委指导", {"size": 40, "color": INK, "bold": True})], align=PP_ALIGN.CENTER)
text(s, Inches(0.7), Inches(4.3), Inches(11.9), Inches(0.8),
     [("海外工程投标及商业调研全流程 AI 智能 Agent 作业系统　｜　海外业务赛道 · 向鑫", {"size": 16, "color": DIM})], align=PP_ALIGN.CENTER)

# ---- 片头视频自动播放：把点击触发(delay=indefinite)改为自动(delay=0) ----
cover = prs.slides[0]
timing = cover.element.find(
    ".//{http://schemas.openxmlformats.org/presentationml/2006/main}timing")
if timing is not None:
    for cond in timing.iter(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}cond"):
        if cond.get("delay") == "indefinite":
            cond.set("delay", "0")

prs.save("competition/out/设计说明-海外工程投标及商业调研全流程AI智能Agent作业系统-向鑫.pptx")
print("PPTX saved")
